"""
Slick Minimal template — python-pptx port.
Signature: thick left green accent bar + thin rule under title.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

W = Inches(10)
H = Inches(5.625)

# Colors
GREEN = RGBColor(0x36, 0x87, 0x27)
GREEN_LIGHT = RGBColor(0xF7, 0xF4, 0xE4)
GREEN_MID = RGBColor(0x1D, 0xE4, 0xCA)
BLUE = RGBColor(0x38, 0x80, 0xF3)
PURPLE = RGBColor(0x5B, 0x2C, 0x8F)
ORANGE = RGBColor(0x04, 0x54, 0x7C)
DARK = RGBColor(0x40, 0x3F, 0x3E)
MID = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xE5, 0xE5, 0xE5)
OFF_WHITE = RGBColor(0xF9, 0xF7, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Calibri"  # Fallback for Fidelity Slab
BODY_FONT = "Calibri"   # Fallback for Fidelity Sans

LM = Inches(0.9)
ACC_W = Inches(0.25)
CW = Inches(8.6)


def _rgb_hex(hex_str):
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, x, y, w, h, text, font_name=BODY_FONT, font_size=12,
                  color=DARK, bold=False, italic=False, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    if valign:
        tf.paragraphs[0].alignment = align

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align

    if line_spacing:
        p.line_spacing = Pt(line_spacing)

    # Set vertical alignment via XML
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}
        bodyPr.set('anchor', anchor_map.get(valign, 't'))

    return txBox


def _accent(slide):
    _add_rect(slide, Inches(0), Inches(0), ACC_W, H, GREEN)


def _slide_title(slide, title, y=None, size=28):
    y = y or Inches(0.4)
    _add_text_box(slide, LM, y, CW, Inches(0.8), title,
                  font_name=TITLE_FONT, font_size=size, color=DARK, bold=True,
                  valign=MSO_ANCHOR.BOTTOM)
    _add_rect(slide, LM, Emu(y + Inches(0.85)), Inches(2.5), Inches(0.04), GREEN)


# ============================================================
# BUILDERS
# ============================================================

def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _accent(slide)
    _add_text_box(slide, LM, Inches(1.0), CW, Inches(1.6), c.get("title", "Title"),
                  TITLE_FONT, 38, DARK, bold=True, valign=MSO_ANCHOR.BOTTOM)
    _add_rect(slide, LM, Inches(2.7), Inches(2.5), Inches(0.04), GREEN)
    if c.get("subtitle"):
        _add_text_box(slide, LM, Inches(2.85), CW, Inches(0.5), c["subtitle"],
                      BODY_FONT, 16, MID)
    meta = []
    if c.get("author"):
        meta.append(c["author"])
    if c.get("date"):
        meta.append(c["date"])
    if meta:
        _add_text_box(slide, LM, Inches(4.2), Inches(5), Inches(0.7),
                      "\n".join(meta), BODY_FONT, 12, MID)


def build_in_brief(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "In Brief"))
    bullets = c.get("bullets", [])
    startY = Inches(1.55)
    rowH = Inches(0.82)
    gap = Inches(0.12)

    for i, b in enumerate(bullets):
        y = Emu(startY + i * (rowH + gap))
        _add_rect(slide, LM, y, Inches(0.06), rowH, GREEN)
        _add_text_box(slide, Emu(LM + Inches(0.25)), y, Emu(CW - Inches(0.25)), rowH,
                      b, BODY_FONT, 13, DARK, valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


def build_section_divider(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Green background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    if c.get("sectionNumber"):
        _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.8),
                      f"0{c['sectionNumber']}", TITLE_FONT, 48, GREEN_MID, bold=True,
                      valign=MSO_ANCHOR.BOTTOM)
    _add_text_box(slide, Inches(0.8), Inches(2.1), Inches(8), Inches(1.0),
                  c.get("title", "Section"), TITLE_FONT, 36, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, Inches(0.8), Inches(3.2), Inches(2.0), Inches(0.04), GREEN_MID)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.8), Inches(3.4), Inches(8), Inches(0.5),
                      c["subtitle"], BODY_FONT, 14, GREEN_LIGHT)


def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Key Metric"), size=20)
    _add_text_box(slide, LM, Inches(1.5), CW, Inches(1.8), c.get("stat", "—"),
                  TITLE_FONT, 80, GREEN, bold=True, align=PP_ALIGN.CENTER,
                  valign=MSO_ANCHOR.MIDDLE)
    if c.get("headline"):
        _add_text_box(slide, Inches(1.5), Inches(3.3), Inches(7), Inches(0.6),
                      c["headline"], BODY_FONT, 16, DARK, bold=True, align=PP_ALIGN.CENTER)
    if c.get("detail"):
        _add_text_box(slide, Inches(1.5), Inches(3.9), Inches(7), Inches(0.7),
                      c["detail"], BODY_FONT, 11, MID, align=PP_ALIGN.CENTER)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(4.85), Inches(9), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True, align=PP_ALIGN.CENTER)


def build_quote(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "In Their Words"), size=20)
    _add_text_box(slide, Emu(LM - Inches(0.1)), Inches(1.4), Inches(0.8), Inches(0.8),
                  "\u201C", "Georgia", 64, GREEN_LIGHT, bold=True)
    _add_text_box(slide, Emu(LM + Inches(0.5)), Inches(1.7), Inches(7.5), Inches(1.8),
                  c.get("quote", ""), BODY_FONT, 17, DARK, italic=True,
                  valign=MSO_ANCHOR.MIDDLE, line_spacing=24)
    _add_rect(slide, Emu(LM + Inches(0.5)), Inches(3.7), Inches(1.5), Inches(0.04), GREEN)
    if c.get("attribution"):
        _add_text_box(slide, Emu(LM + Inches(0.5)), Inches(3.85), Inches(7), Inches(0.35),
                      c["attribution"], BODY_FONT, 12, MID)
    if c.get("context"):
        _add_text_box(slide, Emu(LM + Inches(0.5)), Inches(4.15), Inches(7), Inches(0.35),
                      c["context"], BODY_FONT, 10, MID, italic=True)


def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Comparison"), size=22)

    # Left card
    _add_rect(slide, LM, Inches(1.55), Inches(4.0), Inches(3.5), OFF_WHITE)
    _add_rect(slide, LM, Inches(1.55), Inches(0.06), Inches(3.5), ORANGE)
    _add_text_box(slide, Emu(LM + Inches(0.2)), Inches(1.65), Inches(3.6), Inches(0.35),
                  c.get("leftLabel", "Before"), TITLE_FONT, 14, ORANGE, bold=True)
    for i, item in enumerate(c.get("leftItems", [])):
        _add_text_box(slide, Emu(LM + Inches(0.2)), Emu(Inches(2.1) + i * Inches(0.5)),
                      Inches(3.6), Inches(0.45), item, BODY_FONT, 11, DARK)

    # Divider
    _add_rect(slide, Inches(5.05), Inches(1.75), Inches(0.03), Inches(3.1), LIGHT)

    # Right card
    _add_rect(slide, Inches(5.25), Inches(1.55), Inches(4.25), Inches(3.5), OFF_WHITE)
    _add_rect(slide, Inches(5.25), Inches(1.55), Inches(0.06), Inches(3.5), GREEN)
    _add_text_box(slide, Inches(5.45), Inches(1.65), Inches(3.8), Inches(0.35),
                  c.get("rightLabel", "After"), TITLE_FONT, 14, GREEN, bold=True)
    for i, item in enumerate(c.get("rightItems", [])):
        _add_text_box(slide, Inches(5.45), Emu(Inches(2.1) + i * Inches(0.5)),
                      Inches(3.85), Inches(0.45), item, BODY_FONT, 11, DARK)


def build_text_graph(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Title"), size=24)

    texts = c.get("text", [])
    if not isinstance(texts, list):
        texts = [texts]
    for i, t in enumerate(texts):
        _add_text_box(slide, LM, Emu(Inches(1.6) + i * Inches(1.1)),
                      Inches(4.0), Inches(1.0), t, BODY_FONT, 12, DARK, line_spacing=16)

    _add_rect(slide, Inches(5.1), Inches(1.6), Inches(0.03), Inches(3.3), LIGHT)

    # Add chart
    chart_data_raw = c.get("chartData", [{"name": "S1", "labels": ["A", "B", "C"], "values": [25, 45, 30]}])
    from pptx.chart.data import CategoryChartData
    chart_data = CategoryChartData()
    if chart_data_raw:
        cd = chart_data_raw[0]
        chart_data.categories = cd.get("labels", ["A", "B", "C"])
        chart_data.add_series(cd.get("name", "Series 1"), cd.get("values", [25, 45, 30]))

    ct = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_type = c.get("chartType", "bar")
    if chart_type == "line":
        ct = XL_CHART_TYPE.LINE
    elif chart_type == "pie":
        ct = XL_CHART_TYPE.PIE

    chart = slide.shapes.add_chart(ct, Inches(5.3), Inches(1.3), Inches(4.2), Inches(3.8), chart_data)

    if c.get("note"):
        _add_text_box(slide, Inches(5.3), Inches(5.15), Inches(4.2), Inches(0.3),
                      c["note"], BODY_FONT, 8, MID, italic=True)


def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Process"), size=22)

    steps = c.get("steps", [])
    count = min(len(steps), 5)
    if count == 0:
        return

    total_w = 8.6
    arrow_w = 0.3
    step_w = (total_w - (count - 1) * arrow_w) / count
    startX = 0.9
    stepY = 1.55
    stepH = 3.4

    for i, step in enumerate(steps[:count]):
        x = Inches(startX + i * (step_w + arrow_w))
        y = Inches(stepY)
        sw = Inches(step_w)
        sh = Inches(stepH)

        _add_rect(slide, x, y, sw, sh, OFF_WHITE)
        _add_rect(slide, x, y, Inches(0.06), sh, GREEN)

        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.15)),
                      Inches(0.35), Inches(0.35), str(i + 1),
                      TITLE_FONT, 16, GREEN, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.55)),
                      Emu(sw - Inches(0.3)), Inches(0.55),
                      step.get("title", ""), BODY_FONT, 11, DARK, bold=True)
        if step.get("detail"):
            _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(1.15)),
                          Emu(sw - Inches(0.3)), Inches(1.9),
                          step["detail"], BODY_FONT, 9, MID)

        if i < count - 1:
            _add_text_box(slide, Emu(x + sw + Inches(0.02)), Emu(y + Inches(1.3)),
                          Emu(Inches(arrow_w) - Inches(0.04)), Inches(0.4),
                          "\u2192", BODY_FONT, 14, LIGHT, align=PP_ALIGN.CENTER,
                          valign=MSO_ANCHOR.MIDDLE)


def build_matrix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Framework"), size=22)

    quads = c.get("quadrants", [{}, {}, {}, {}])
    accents = [GREEN, BLUE, PURPLE, ORANGE]
    qW = Inches(4.05)
    qH = Inches(1.7)
    gap = Inches(0.2)
    sX = LM
    sY = Inches(1.55)

    for i, q in enumerate(quads[:4]):
        col = i % 2
        row = i // 2
        x = Emu(sX + col * (qW + gap))
        y = Emu(sY + row * (qH + gap))

        _add_rect(slide, x, y, qW, qH, OFF_WHITE)
        _add_rect(slide, x, y, Inches(0.06), qH, accents[i])
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.1)),
                      Emu(qW - Inches(0.4)), Inches(0.3),
                      q.get("label", ""), TITLE_FONT, 11, accents[i], bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.45)),
                      Emu(qW - Inches(0.4)), Emu(qH - Inches(0.6)),
                      q.get("detail", ""), BODY_FONT, 10, DARK)


def build_methods(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Approach"), size=22)

    for i, f in enumerate(c.get("fields", [])):
        y = Emu(Inches(1.55) + i * Inches(0.72))
        _add_rect(slide, LM, y, Inches(0.06), Inches(0.55), GREEN)
        _add_text_box(slide, Emu(LM + Inches(0.2)), y, Inches(1.8), Inches(0.55),
                      f.get("label", ""), BODY_FONT, 12, GREEN, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Emu(LM + Inches(2.1)), y, Inches(6.4), Inches(0.55),
                      f.get("value", ""), BODY_FONT, 12, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_hypotheses(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Hypotheses"), size=22)

    for i, h in enumerate(c.get("hypotheses", [])):
        y = Emu(Inches(1.55) + i * Inches(0.7))
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        _add_rect(slide, LM, y, Inches(8.6), Inches(0.58), bg)
        _add_rect(slide, LM, y, Inches(0.06), Inches(0.58), GREEN)

        _add_text_box(slide, Emu(LM + Inches(0.15)), y, Inches(0.45), Inches(0.58),
                      f"H{i + 1}", TITLE_FONT, 12, GREEN, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Emu(LM + Inches(0.65)), y, Inches(6.3), Inches(0.58),
                      h.get("text", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        if h.get("status"):
            sc = GREEN if h["status"] == "Confirmed" else (ORANGE if h["status"] == "Rejected" else MID)
            _add_text_box(slide, Inches(7.8), Emu(y + Inches(0.14)), Inches(1.0), Inches(0.3),
                          h["status"], BODY_FONT, 9, sc, bold=True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def build_wsn_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Key Finding"), size=22)

    cols = [
        ("What", GREEN, c.get("what", {})),
        ("So What", BLUE, c.get("soWhat", {})),
        ("Now What", PURPLE, c.get("nowWhat", {})),
    ]
    colW = Inches(2.75)
    gap = Inches(0.2)
    startX = LM
    startY = Inches(1.55)
    cardH = Inches(3.5)

    for i, (label, color, data) in enumerate(cols):
        x = Emu(startX + i * (colW + gap))
        _add_rect(slide, x, startY, colW, cardH, OFF_WHITE)
        _add_rect(slide, x, startY, Inches(0.06), cardH, color)

        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(startY + Inches(0.15)),
                      Emu(colW - Inches(0.4)), Inches(0.35),
                      label, TITLE_FONT, 13, color, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(startY + Inches(0.55)),
                      Emu(colW - Inches(0.4)), Inches(1.1),
                      data.get("headline", ""), BODY_FONT, 11, DARK, bold=True)
        if data.get("detail"):
            _add_text_box(slide, Emu(x + Inches(0.2)), Emu(startY + Inches(1.7)),
                          Emu(colW - Inches(0.4)), Inches(1.6),
                          data["detail"], BODY_FONT, 9.5, MID)


def build_wsn_reveal(prs, c):
    """Builds 3 progressive slides for What → So What → Now What."""
    def _draw_title(slide, title):
        _accent(slide)
        _add_text_box(slide, LM, Inches(0.4), CW, Inches(0.7), title,
                      TITLE_FONT, 24, DARK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, LM, Inches(1.15), Inches(2.5), Inches(0.04), GREEN)

    def _draw_zone(slide, x, label, color, data, condensed=False):
        h = Inches(1.95) if condensed else Inches(3.3)
        y = Inches(1.55)
        w = Inches(4.05)

        _add_rect(slide, x, y, w, h, OFF_WHITE)
        _add_rect(slide, x, y, Inches(0.06), h, color)

        fs_label = 10 if condensed else 11
        fs_head = 10 if condensed else 12
        fs_detail = 8.5 if condensed else 10

        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.1)),
                      Inches(2), Inches(0.3), label, TITLE_FONT, fs_label, color, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.45)),
                      Emu(w - Inches(0.4)), Inches(0.55) if condensed else Inches(0.8),
                      data.get("headline", ""), BODY_FONT, fs_head, DARK, bold=True)
        if data.get("detail"):
            detail_y = Emu(y + Inches(1.0)) if condensed else Emu(y + Inches(1.3))
            detail_h = Inches(0.7) if condensed else Inches(1.6)
            _add_text_box(slide, Emu(x + Inches(0.2)), detail_y,
                          Emu(w - Inches(0.4)), detail_h,
                          data["detail"], BODY_FONT, fs_detail, MID)

    leftX = LM
    rightX = Inches(5.2)

    # Slide 1: What only
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_title(s1, c.get("title", "Key Finding"))
    _draw_zone(s1, leftX, "What We Found", GREEN, c.get("what", {}))

    # Slide 2: What + So What
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_title(s2, c.get("title", "Key Finding"))
    _draw_zone(s2, leftX, "What We Found", GREEN, c.get("what", {}))
    _draw_zone(s2, rightX, "So What", BLUE, c.get("soWhat", {}))

    # Slide 3: All three condensed + Now What card
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_title(s3, c.get("title", "Key Finding"))
    _draw_zone(s3, leftX, "What We Found", GREEN, c.get("what", {}), condensed=True)
    _draw_zone(s3, rightX, "So What", BLUE, c.get("soWhat", {}), condensed=True)

    # Now What full-width card
    nwY = Inches(3.7)
    nwH = Inches(1.6)
    _add_rect(s3, LM, nwY, CW, nwH, OFF_WHITE)
    _add_rect(s3, LM, nwY, Inches(0.06), nwH, PURPLE)
    _add_text_box(s3, Emu(LM + Inches(0.2)), Emu(nwY + Inches(0.1)),
                  Inches(2), Inches(0.25), "Now What", TITLE_FONT, 10, PURPLE, bold=True)
    nw = c.get("nowWhat", {})
    _add_text_box(s3, Emu(LM + Inches(0.2)), Emu(nwY + Inches(0.4)),
                  Emu(CW - Inches(0.4)), Inches(0.6),
                  nw.get("headline", ""), BODY_FONT, 16, DARK, bold=True)
    if nw.get("detail"):
        _add_text_box(s3, Emu(LM + Inches(0.2)), Emu(nwY + Inches(1.0)),
                      Emu(CW - Inches(0.4)), Inches(0.45),
                      nw["detail"], BODY_FONT, 10, MID)


def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Findings & Recommendations"), size=22)

    items = c.get("items", [])
    startY = Inches(1.6)
    rowH = Inches(0.72)
    gap = Inches(0.14)

    for i, item in enumerate(items[:5]):
        y = Emu(startY + i * (rowH + gap))
        # Finding card
        _add_rect(slide, LM, y, Inches(3.9), rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.06), rowH, GREEN)
        _add_text_box(slide, Emu(LM + Inches(0.2)), y, Inches(3.5), rowH,
                      item.get("finding", ""), BODY_FONT, 10.5, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        # Arrow
        _add_text_box(slide, Inches(4.95), y, Inches(0.35), rowH,
                      "\u2192", BODY_FONT, 14, GREEN, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)
        # Rec card
        _add_rect(slide, Inches(5.4), y, Inches(4.1), rowH, OFF_WHITE)
        _add_rect(slide, Inches(5.4), y, Inches(0.06), rowH, BLUE)
        _add_text_box(slide, Inches(5.6), y, Inches(3.75), rowH,
                      item.get("recommendation", ""), BODY_FONT, 10.5, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Complete Findings"), size=20)

    items = c.get("items", [])
    startY = Inches(1.5)
    rowH = Inches(0.44)
    gap = Inches(0.06)

    for i, item in enumerate(items[:8]):
        y = Emu(startY + i * (rowH + gap))
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        _add_rect(slide, LM, y, Inches(4.1), rowH, bg)
        _add_rect(slide, LM, y, Inches(0.04), rowH, GREEN)
        _add_text_box(slide, Emu(LM + Inches(0.15)), y, Inches(3.9), rowH,
                      item.get("finding", ""), BODY_FONT, 9, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, Inches(5.15), y, Inches(4.35), rowH, bg)
        _add_rect(slide, Inches(5.15), y, Inches(0.04), rowH, BLUE)
        _add_text_box(slide, Inches(5.3), y, Inches(4.1), rowH,
                      item.get("recommendation", ""), BODY_FONT, 9, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Open Questions"), size=26)

    questions = c.get("questions", [])
    cardW = Inches(4.1)
    cardH = Inches(1.55)
    gX = Inches(0.4)
    gY = Inches(0.25)
    gridX = LM
    gridY = Inches(1.55)

    for i, question in enumerate(questions[:4]):
        col = i % 2
        row = i // 2
        x = Emu(gridX + col * (cardW + gX))
        y = Emu(gridY + row * (cardH + gY))

        _add_rect(slide, x, y, cardW, cardH, OFF_WHITE)
        _add_rect(slide, x, y, Inches(0.06), cardH, GREEN)

        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.1)),
                      Inches(0.4), Inches(0.4), str(i + 1),
                      TITLE_FONT, 22, GREEN, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.55)),
                      Emu(cardW - Inches(0.4)), Inches(0.85),
                      question, BODY_FONT, 12, DARK)


def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Agenda"), size=26)

    items = c.get("items", [])
    for i, item in enumerate(items):
        y = Emu(Inches(1.55) + i * Inches(0.7))
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        _add_rect(slide, LM, y, CW, Inches(0.58), bg)
        _add_rect(slide, LM, y, Inches(0.06), Inches(0.58), GREEN)

        _add_text_box(slide, Emu(LM + Inches(0.15)), y, Inches(0.4), Inches(0.58),
                      str(i + 1), TITLE_FONT, 16, GREEN, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        title_text = item if isinstance(item, str) else item.get("title", "")
        _add_text_box(slide, Emu(LM + Inches(0.65)), y, Inches(5.5), Inches(0.58),
                      title_text, BODY_FONT, 14, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)

        if isinstance(item, dict) and item.get("detail"):
            _add_text_box(slide, Inches(7.5), y, Inches(1.8), Inches(0.58),
                          item["detail"], BODY_FONT, 10, MID,
                          align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_progressive_reveal(prs, c):
    takeaways = c.get("takeaways", [])
    for n in range(min(len(takeaways), 5)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _accent(slide)
        _slide_title(slide, c.get("title", "Building the Picture"))

        cur = takeaways[n]
        _add_rect(slide, LM, Inches(1.55), CW, Inches(2.2), OFF_WHITE)
        _add_rect(slide, LM, Inches(1.55), Inches(0.06), Inches(2.2), GREEN)

        _add_text_box(slide, Emu(LM + Inches(0.2)), Inches(1.65),
                      Emu(CW - Inches(0.4)), Inches(0.6),
                      cur.get("headline", ""), BODY_FONT, 15, DARK, bold=True)
        if cur.get("detail"):
            _add_text_box(slide, Emu(LM + Inches(0.2)), Inches(2.3),
                          Emu(CW - Inches(0.4)), Inches(1.2),
                          cur["detail"], BODY_FONT, 11, MID)

        # Running takeaway bar
        _add_rect(slide, Inches(0), Inches(3.95), W, Inches(0.04), GREEN)
        _add_text_box(slide, LM, Inches(4.05), Inches(3), Inches(0.25),
                      "Running Takeaways", TITLE_FONT, 9, GREEN, bold=True)

        for j in range(n + 1):
            ty = Emu(Inches(4.35) + j * Inches(0.3))
            active = j == n
            _add_rect(slide, LM, Emu(ty + Inches(0.04)), Inches(0.12), Inches(0.12), GREEN)
            _add_text_box(slide, Emu(LM + Inches(0.25)), ty,
                          Emu(CW - Inches(0.25)), Inches(0.28),
                          takeaways[j].get("summary", takeaways[j].get("headline", "")),
                          BODY_FONT, 9, DARK if active else MID, bold=active,
                          valign=MSO_ANCHOR.MIDDLE)


def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    _add_text_box(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(1.2),
                  c.get("title", "Thank You"), TITLE_FONT, 44, WHITE, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
    _add_rect(slide, Inches(3.75), Inches(2.75), Inches(2.5), Inches(0.04), GREEN_MID)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.5), Inches(2.95), Inches(9), Inches(0.5),
                      c["subtitle"], BODY_FONT, 16, WHITE, align=PP_ALIGN.CENTER)
    if c.get("contact"):
        _add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
                      c["contact"], BODY_FONT, 12, GREEN_LIGHT, align=PP_ALIGN.CENTER)


# ============================================================
# DISPATCH
# ============================================================
BUILDERS = {
    "title": build_title,
    "in_brief": build_in_brief,
    "section_divider": build_section_divider,
    "stat_callout": build_stat_callout,
    "quote": build_quote,
    "comparison": build_comparison,
    "text_graph": build_text_graph,
    "process_flow": build_process_flow,
    "matrix": build_matrix,
    "methods": build_methods,
    "hypotheses": build_hypotheses,
    "wsn_dense": build_wsn_dense,
    "wsn_reveal": build_wsn_reveal,
    "findings_recs": build_findings_recs,
    "findings_recs_dense": build_findings_recs_dense,
    "open_questions": build_open_questions,
    "agenda": build_agenda,
    "progressive_reveal": build_progressive_reveal,
    "closer": build_closer,
}


def build_deck(slide_configs, output_path):
    """Build a complete deck from a list of (slide_type, data_dict) tuples."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_type, data in slide_configs:
        if slide_type == "skip":
            continue
        builder = BUILDERS.get(slide_type)
        if builder:
            builder(prs, data)

    prs.save(output_path)
    return output_path
