"""
Slide-type auto-detection engine.
Reads a .pptx, extracts text structure from each slide,
and scores every template type to return ranked candidates.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import re

SLIDE_TYPES = [
    "title", "agenda", "in_brief", "section_divider", "stat_callout",
    "quote", "comparison", "text_graph", "process_flow", "matrix",
    "methods", "hypotheses", "wsn_dense", "wsn_reveal",
    "findings_recs", "findings_recs_dense", "open_questions",
    "progressive_reveal", "closer", "skip"
]

SLIDE_TYPE_LABELS = {
    "title": "Title Slide",
    "agenda": "Agenda",
    "in_brief": "In Brief (Bullets)",
    "section_divider": "Section Divider",
    "stat_callout": "Stat Callout",
    "quote": "Quote",
    "comparison": "Comparison (Left/Right)",
    "text_graph": "Text + Graph",
    "process_flow": "Process Flow (Steps)",
    "matrix": "2×2 Matrix",
    "methods": "Methods / Key-Value Fields",
    "hypotheses": "Hypotheses",
    "wsn_dense": "What / So What / Now What (Dense)",
    "wsn_reveal": "What / So What / Now What (Reveal)",
    "findings_recs": "Findings & Recommendations",
    "findings_recs_dense": "Findings & Recs (Dense)",
    "open_questions": "Open Questions (2×2 Grid)",
    "progressive_reveal": "Progressive Reveal",
    "closer": "Closer / Thank You",
    "skip": "Skip (exclude)",
}

SLIDE_TYPE_DESCRIPTIONS = {
    "title": "Big title, subtitle, author/date. Use for the opening slide.",
    "agenda": "Numbered list of topics with optional timing. Sets the roadmap.",
    "in_brief": "3–5 key bullet points with accent bars. Your workhorse summary slide.",
    "section_divider": "Bold text on colored background. Marks a new section.",
    "stat_callout": "One huge number/percentage with headline. Maximum impact.",
    "quote": "Large italic quote with attribution. Puts a human voice on the data.",
    "comparison": "Side-by-side cards (before/after, old/new). Shows contrast.",
    "text_graph": "Text on left, chart on right. Pairs narrative with data.",
    "process_flow": "3–5 sequential step cards with arrows. Shows a journey.",
    "matrix": "2×2 grid of quadrants. Frameworks and tradeoffs.",
    "methods": "Label:value pairs stacked vertically. Methodology or specs.",
    "hypotheses": "Numbered hypotheses with confirmed/rejected status badges.",
    "wsn_dense": "Three-column What/So What/Now What. Dense single-slide insight.",
    "wsn_reveal": "Builds across 3 slides: What → So What → Now What.",
    "findings_recs": "Paired finding→recommendation cards (up to 5).",
    "findings_recs_dense": "Compact paired rows (up to 8). Lots of findings.",
    "open_questions": "2×2 grid of question cards. Invites discussion.",
    "progressive_reveal": "Multi-slide build with running takeaway bar at bottom.",
    "closer": "Thank you slide on green background. Contact info optional.",
    "skip": "Exclude this slide from the output entirely.",
}


def extract_slides(pptx_path):
    """Extract text structure from each slide."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        sd = {
            "index": i, "number": i + 1, "shapes": [], "all_text": [],
            "text_boxes": [], "has_chart": False, "has_table": False,
            "has_image": False, "shape_count": 0,
        }
        for shape in slide.shapes:
            sd["shape_count"] += 1
            if shape.has_chart:
                sd["has_chart"] = True
            if shape.has_table:
                sd["has_table"] = True
            if hasattr(shape, "image"):
                sd["has_image"] = True
            if shape.has_text_frame:
                paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                if paragraphs:
                    max_font = 0
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt > max_font:
                                max_font = run.font.size.pt
                        # Also check paragraph-level font (some decks set it there)
                        if hasattr(para, 'font') and para.font.size and para.font.size.pt > max_font:
                            max_font = para.font.size.pt
                    # Fallback: estimate from shape height and text count
                    if max_font == 0 and shape.height and len(paragraphs) > 0:
                        est = round(shape.height.pt / max(len(paragraphs) * 1.5, 1))
                        if est > 6:
                            max_font = est
                    sd["text_boxes"].append({
                        "type": "text", "paragraphs": paragraphs,
                        "text": "\n".join(paragraphs), "para_count": len(paragraphs),
                        "max_font_size": max_font, "left": shape.left, "top": shape.top,
                        "width": shape.width,
                    })
                    sd["all_text"].extend(paragraphs)
        sd["total_text"] = "\n".join(sd["all_text"])
        sd["total_words"] = len(sd["total_text"].split())
        sd["text_box_count"] = len(sd["text_boxes"])
        slides.append(sd)
    return slides


def _score_all_types(slide):
    """Score every slide type. Returns sorted list of (type, score, reason)."""
    texts = slide["all_text"]
    boxes = slide["text_boxes"]
    total_words = slide["total_words"]
    full = slide["total_text"].lower()
    box_count = slide["text_box_count"]
    max_font = max((b["max_font_size"] for b in boxes), default=0)
    bullet_like = [t for t in texts if len(t) > 15]

    if not texts:
        return [("skip", 0.5, "Empty slide")]

    scores = []

    # ── TITLE ──
    s, r = 0.0, []
    if slide["number"] == 1:
        s += 0.5; r.append("first slide")
    if total_words < 40:
        s += 0.2; r.append(f"{total_words} words")
    if max_font >= 30:
        s += 0.2; r.append(f"{max_font}pt font")
    if total_words < 15 and max_font >= 36:
        s += 0.15
    scores.append(("title", min(s, 0.95), "; ".join(r) or "no strong signals"))

    # ── CLOSER ──
    s, r = 0.0, []
    closer_kw = ["thank you", "thanks", "q&a", "contact"]
    m = [k for k in closer_kw if k in full]
    if m: s += 0.5; r.append(f"keywords: {', '.join(m)}")
    if total_words < 30: s += 0.2; r.append("short text")
    if slide["number"] > 3: s += 0.1
    scores.append(("closer", min(s, 0.95), "; ".join(r) or "no closing keywords"))

    # ── SECTION DIVIDER ──
    s, r = 0.0, []
    if total_words < 15: s += 0.3; r.append(f"{total_words} words")
    if max_font >= 24: s += 0.25; r.append(f"{max_font}pt font")
    if max_font >= 32: s += 0.15
    if slide["number"] > 1: s += 0.05
    scores.append(("section_divider", min(s, 0.95), "; ".join(r) or "no strong signals"))

    # ── AGENDA ──
    s, r = 0.0, []
    agenda_kw = ["agenda", "outline", "overview", "today's plan", "topics"]
    m = [k for k in agenda_kw if k in full]
    if m: s += 0.55; r.append(f"keywords: {', '.join(m)}")
    bc = len([t for t in texts if len(t) > 8])
    if 3 <= bc <= 8 and total_words < 100: s += 0.2; r.append(f"{bc} items")
    scores.append(("agenda", min(s, 0.95), "; ".join(r) or "no agenda keywords"))

    # ── IN BRIEF ──
    s, r = 0.0, []
    if len(bullet_like) >= 3: s += 0.45; r.append(f"{len(bullet_like)} bullet-length items")
    if len(bullet_like) >= 5: s += 0.1
    if total_words > 40: s += 0.1; r.append("substantial text")
    scores.append(("in_brief", min(s, 0.95), "; ".join(r) or "few bullet-length items"))

    # ── STAT CALLOUT ──
    s, r = 0.0, []
    big_num = re.compile(r'^\s*[\d,.]+[%×xX]?\s*$')
    for box in boxes:
        if box["para_count"] == 1 and big_num.match(box["text"]):
            s += 0.5; r.append(f"standalone number: {box['text'].strip()}")
            if box["max_font_size"] >= 36: s += 0.2; r.append("large font")
            break
    nums = re.findall(r'\b\d+[%×xX]\b', full)
    if nums and s < 0.3: s += 0.3; r.append(f"numbers: {', '.join(nums[:3])}")
    if total_words < 30: s += 0.15
    scores.append(("stat_callout", min(s, 0.95), "; ".join(r) or "no big numbers"))

    # ── QUOTE ──
    s, r = 0.0, []
    has_q = "\u201c" in full or "\u201d" in full or full.count('"') >= 2
    if has_q: s += 0.35; r.append("quotation marks")
    has_attr = any(w in full for w in ["—", "\u2014", "- ", "attributed", "said"])
    if has_attr: s += 0.25; r.append("attribution pattern")
    if has_q and total_words < 80: s += 0.1
    scores.append(("quote", min(s, 0.95), "; ".join(r) or "no quoted text"))

    # ── OPEN QUESTIONS ──
    s, r = 0.0, []
    qm = sum(1 for t in texts if t.strip().endswith("?"))
    if qm >= 3: s += 0.7; r.append(f"{qm} questions")
    elif qm >= 2: s += 0.45; r.append(f"{qm} questions")
    elif qm == 1: s += 0.15; r.append("1 question")
    scores.append(("open_questions", min(s, 0.95), "; ".join(r) or "no questions"))

    # ── HYPOTHESES ──
    s, r = 0.0, []
    hyp_kw = ["hypothesis", "hypotheses", "h1:", "h2:", "h3:"]
    m = [k for k in hyp_kw if k in full]
    if m: s += 0.5; r.append(f"keywords: {', '.join(m)}")
    st_words = ["confirmed", "rejected", "partial", "supported", "not supported"]
    st_m = [w for w in st_words if w in full]
    if st_m: s += 0.25; r.append(f"status: {', '.join(st_m)}")
    scores.append(("hypotheses", min(s, 0.95), "; ".join(r) or "no hypothesis keywords"))

    # ── WSN DENSE / REVEAL ──
    s, r = 0.0, []
    wsn_kw = ["what", "so what", "now what"]
    wc = sum(1 for k in wsn_kw if k in full)
    if wc >= 3: s += 0.75; r.append("all three WSN sections")
    elif wc >= 2: s += 0.55; r.append(f"{wc}/3 WSN sections")
    scores.append(("wsn_dense", min(s, 0.95), "; ".join(r) or "no WSN structure"))
    scores.append(("wsn_reveal", min(s * 0.9, 0.95), ("; ".join(r) + " (3-slide build)") if r else "no WSN structure"))

    # ── COMPARISON ──
    s, r = 0.0, []
    if any(k in full for k in ["vs.", "vs ", "versus"]): s += 0.5; r.append("contains 'vs'")
    comp_kw = ["before", "after", "traditional", "current", "new", "old"]
    cm = [k for k in comp_kw if k in full]
    if len(cm) >= 2: s += 0.2; r.append(f"words: {', '.join(cm)}")
    if box_count >= 2:
        lefts = sorted(set(b["left"] for b in boxes if b["left"] is not None))
        if len(lefts) >= 2 and (lefts[-1] - lefts[0]) > Emu(3000000):
            s += 0.2; r.append("two-column layout")
    scores.append(("comparison", min(s, 0.95), "; ".join(r) or "no comparison signals"))

    # ── METHODS ──
    s, r = 0.0, []
    m_kw = ["method", "approach", "sample", "design", "analysis", "measure", "participant", "limitation"]
    mm = [k for k in m_kw if k in full]
    if len(mm) >= 3: s += 0.55; r.append(f"{len(mm)} keywords: {', '.join(mm)}")
    elif len(mm) >= 2: s += 0.3; r.append(f"{len(mm)} keywords: {', '.join(mm)}")
    kv = sum(1 for t in texts if ":" in t and len(t) > 10)
    if kv >= 3: s += 0.2; r.append(f"{kv} key:value pairs")
    scores.append(("methods", min(s, 0.95), "; ".join(r) or "no methodology keywords"))

    # ── FINDINGS & RECS ──
    s, r = 0.0, []
    fr_kw = ["finding", "recommendation", "implication", "action", "suggest"]
    fm = [k for k in fr_kw if k in full]
    if len(fm) >= 2: s += 0.4; r.append(f"keywords: {', '.join(fm)}")
    ac = full.count("\u2192") + full.count("->") + full.count("\u279c")
    if ac >= 2: s += 0.35; r.append(f"{ac} arrow patterns")
    bc2 = len([t for t in texts if len(t) > 10])
    dense = bc2 > 6
    scores.append(("findings_recs", min(s if not dense else s * 0.7, 0.95),
                   "; ".join(r) or "no finding/rec patterns"))
    scores.append(("findings_recs_dense", min(s * (1.1 if dense else 0.7), 0.95),
                   ("; ".join(r) + (f"; {bc2} items" if dense else "")) or "no finding/rec patterns"))

    # ── PROCESS FLOW ──
    s, r = 0.0, []
    sk = ["step 1", "step 2", "phase 1", "phase 2", "stage 1", "first,", "then,", "finally,"]
    sm = [k for k in sk if k in full]
    if len(sm) >= 2: s += 0.5; r.append(f"keywords: {', '.join(sm)}")
    nd = re.findall(r'(?:^|\n)\s*\d+[\.\)]\s', slide["total_text"])
    if len(nd) >= 3: s += 0.3; r.append(f"{len(nd)} numbered items")
    scores.append(("process_flow", min(s, 0.95), "; ".join(r) or "no step patterns"))

    # ── MATRIX ──
    s, r = 0.0, []
    mx = ["quadrant", "matrix", "framework", "2x2", "2×2", "high/low"]
    mm2 = [k for k in mx if k in full]
    if mm2: s += 0.55; r.append(f"keywords: {', '.join(mm2)}")
    scores.append(("matrix", min(s, 0.95), "; ".join(r) or "no matrix keywords"))

    # ── TEXT + GRAPH ──
    s, r = 0.0, []
    if slide["has_chart"]: s += 0.8; r.append("contains a chart")
    scores.append(("text_graph", min(s, 0.95), "; ".join(r) or "no chart"))

    # ── PROGRESSIVE REVEAL ──
    s, r = 0.0, []
    if len(bullet_like) >= 3 and total_words > 80:
        s += 0.2; r.append("multi-point content")
    scores.append(("progressive_reveal", min(s, 0.95), "; ".join(r) or "better as single slide"))

    scores.sort(key=lambda x: x[1], reverse=True)
    return scores


def detect_slide_candidates(slide, top_n=3):
    """Return the top N candidates with scores and reasons."""
    scores = _score_all_types(slide)
    meaningful = [(t, s, r) for t, s, r in scores if s > 0.05]
    if not meaningful:
        return [("in_brief", 0.3, "fallback")]
    return meaningful[:top_n]


def analyze_deck(pptx_path):
    """Full analysis: extract slides, detect types, return ranked candidates."""
    slides = extract_slides(pptx_path)
    results = []
    for slide in slides:
        candidates = detect_slide_candidates(slide, top_n=3)
        best_type, best_conf, best_reason = candidates[0]
        preview = slide["total_text"][:120].replace("\n", " ")
        if len(slide["total_text"]) > 120:
            preview += "..."
        results.append({
            "number": slide["number"],
            "detected_type": best_type,
            "confidence": best_conf,
            "reason": best_reason,
            "candidates": [
                {"type": t, "confidence": round(c, 2), "reason": r}
                for t, c, r in candidates
            ],
            "preview": preview,
            "total_words": slide["total_words"],
            "text_boxes": len(slide["text_boxes"]),
            "has_chart": slide["has_chart"],
            "has_table": slide["has_table"],
            "has_image": slide["has_image"],
            "all_text": slide["all_text"],
            "raw_boxes": slide["text_boxes"],
        })
    return results
